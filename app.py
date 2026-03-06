import streamlit as st
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

st.set_page_config(page_title="Poster Frames PPT", layout="wide")
st.title("📊 Poster Frames POP PPT Generator")

# -------------------------
# Google Drive Authentication
# -------------------------
def authenticate_drive():
    credentials = service_account.Credentials.from_service_account_info(
        st.secrets["gdrive"],
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )
    return build("drive", "v3", credentials=credentials)


def extract_folder_id(link):
    if "folders/" not in link:
        st.error("Please provide valid Google Drive folder link")
        st.stop()
    return link.split("folders/")[1].split("?")[0]


def get_subfolders(service, parent_id):
    query = f"'{parent_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
    results = service.files().list(q=query, fields="files(id, name)").execute()
    return results.get("files", [])


def get_images_in_folder(service, folder_id):
    query = f"'{folder_id}' in parents and mimeType contains 'image/' and trashed=false"
    results = service.files().list(q=query, fields="files(id, name)").execute()
    return results.get("files", [])


def download_image(service, file_id):
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)

    done = False
    while not done:
        _, done = downloader.next_chunk()

    fh.seek(0)
    return fh


# -------------------------
# User Inputs
# -------------------------
campaign_input = st.text_input("📌 Campaign Name")
drive_link = st.text_input("🔗 Google Drive Folder Link")
generate_btn = st.button("🚀 Generate Presentation")


if generate_btn:

    if not campaign_input or not drive_link:
        st.warning("Please fill all fields")
        st.stop()

    try:
        service = authenticate_drive()
        main_folder_id = extract_folder_id(drive_link)
        subfolders = get_subfolders(service, main_folder_id)

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        TEAL = RGBColor(0, 140, 170)

        # -------------------------
        # Layout for 2 Images
        # -------------------------
        image_width = Inches(3)
        gap = Inches(0.4)

        left_positions = [
            Inches(0.6),
            Inches(5.0)
        ]

        top_position = Inches(2.2)

        # -------------------------
        # Loop Through Folders
        # -------------------------
        for folder in subfolders:

            images = get_images_in_folder(service, folder["id"])
            if not images:
                continue

            for i in range(0, len(images), 2):

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # -------------------------
                # Header
                # -------------------------
                header = slide.shapes.add_shape(
                    1, Inches(0.2), Inches(0.2), Inches(4.5), Inches(0.7)
                )

                header.fill.solid()
                header.fill.fore_color.rgb = TEAL
                header.line.fill.background()

                header_tf = header.text_frame
                header_tf.clear()

                hp = header_tf.paragraphs[0]
                hp.text = folder["name"]
                hp.font.size = Pt(18)
                hp.font.bold = True
                hp.font.color.rgb = RGBColor(255, 255, 255)
                hp.alignment = PP_ALIGN.CENTER

                # -------------------------
                # Campaign Name
                # -------------------------
                campaign_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.1), Inches(8), Inches(0.8)
                )

                campaign_tf = campaign_box.text_frame
                campaign_tf.clear()

                cp = campaign_tf.paragraphs[0]
                cp.text = f"Campaign Name: {campaign_input}"
                cp.font.size = Pt(26)
                cp.font.bold = True
                cp.font.color.rgb = RGBColor(0, 0, 0)
                cp.alignment = PP_ALIGN.LEFT

                # -------------------------
                # Add Images (2 per slide)
                # -------------------------
                slide_images = images[i:i+2]

                for idx, img in enumerate(slide_images):

                    img_stream = download_image(service, img["id"])

                    picture = slide.shapes.add_picture(
                        img_stream,
                        left_positions[idx],
                        top_position,
                        width=image_width
                    )

                    # Border
                    border = slide.shapes.add_shape(
                        1,
                        picture.left,
                        picture.top,
                        picture.width,
                        picture.height
                    )

                    border.fill.background()
                    border.line.color.rgb = RGBColor(0, 0, 0)
                    border.line.width = Pt(1.5)

        # -------------------------
        # Save PPT
        # -------------------------
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        st.success("Presentation generated successfully!")

        st.download_button(
            label="📥 Download PPT",
            data=ppt_io,
            file_name=f"{campaign_input}_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    except Exception as e:
        st.error(f"Error: {e}")
